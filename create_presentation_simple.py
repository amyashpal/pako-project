"""
Simplified script to create PowerPoint presentation for Gujarat Crop Recommendation System
Run: python create_presentation_simple.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide_with_title_and_content(prs, title_text, content_items, colors):
    """Helper function to add a slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
    
    # Add title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = colors['primary']
    
    # Add content box
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Determine if it's a main point or sub-point
        if item.startswith('  ') or item.startswith('   '):
            p.text = item.strip()
            p.level = 1
            p.font.size = Pt(16)
        elif item == '':
            continue
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            if item.endswith(':'):
                p.font.bold = True
    
    return slide

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    colors = {
        'primary': RGBColor(76, 175, 80),  # Green
        'secondary': RGBColor(33, 150, 243),  # Blue
        'text': RGBColor(33, 33, 33)  # Dark gray
    }
    
    print("Creating presentation...")
    
    # Slide 1: Title Slide
    print("  Adding slide 1: Title")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Gujarat Crop Recommendation System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['primary']
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Crop Recommendation for Farmers"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = colors['text']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement
    print("  Adding slide 2: Problem Statement")
    content = [
        "Challenges Faced by Farmers:",
        "",
        "• Uncertainty in crop selection based on soil conditions",
        "• Lack of scientific guidance for crop planning",
        "• Low productivity due to unsuitable crop choices",
        "• Limited access to agricultural expertise",
        "• Language barriers in accessing technology",
        "• Risk of crop failure and financial losses"
    ]
    add_slide_with_title_and_content(prs, "Problem Statement", content, colors)
    
    # Slide 3: Solution Overview
    print("  Adding slide 3: Solution")
    content = [
        "AI-Powered Crop Recommendation System",
        "",
        "✓ Smart crop recommendations based on soil analysis",
        "✓ Yield prediction using machine learning",
        "✓ Treatment plans (fertilizers, pesticides, irrigation)",
        "✓ Multi-language support (12+ Indian languages)",
        "✓ User-friendly web interface for farmers",
        "✓ Prediction history and analytics"
    ]
    add_slide_with_title_and_content(prs, "Our Solution", content, colors)
    
    # Slide 4: Technology Stack
    print("  Adding slide 4: Technology Stack")
    content = [
        "Technology Stack:",
        "",
        "Frontend: Streamlit, HTML/CSS, JavaScript",
        "Backend: Python, PostgreSQL, JWT Authentication",
        "Machine Learning: XGBoost, Random Forest, Scikit-learn",
        "Data Processing: Pandas, NumPy, SMOTE",
        "Visualization: Matplotlib, Seaborn, SHAP",
        "Translation: Deep Translator (Google Translate API)",
        "Deployment: Docker, Railway Platform"
    ]
    add_slide_with_title_and_content(prs, "Technology Stack", content, colors)
    
    # Slide 5: Machine Learning Models
    print("  Adding slide 5: ML Models")
    content = [
        "Three Powerful ML Models:",
        "",
        "1. XGBoost Classifier (Crop Recommendation)",
        "   • Multi-class classification for 21 crops",
        "   • 83.16% accuracy, 99.07% top-3 accuracy",
        "",
        "2. Random Forest (Suitability Prediction)",
        "   • Binary classification (Grow/Not Grow)",
        "   • 88.92% accuracy, 91.52% ROC AUC",
        "",
        "3. XGBoost Regressor (Yield Prediction)",
        "   • Predicts expected crop yield",
        "   • 95.96% R² score, RMSE: 11.52"
    ]
    add_slide_with_title_and_content(prs, "Machine Learning Models", content, colors)
    
    # Slide 6: Key Features
    print("  Adding slide 6: Key Features")
    content = [
        "System Features:",
        "",
        "🌾 Normal Mode: Top 3-4 crop recommendations",
        "🎯 Advanced Mode: Specific crop suitability analysis",
        "📊 Yield Prediction: Expected crop yield estimation",
        "💊 Treatment Plans: Fertilizers, pesticides, irrigation",
        "🌍 Multi-language: 12+ Indian languages supported",
        "📜 History: Track all predictions and analytics",
        "🔐 Secure: Encrypted passwords, JWT authentication"
    ]
    add_slide_with_title_and_content(prs, "Key Features", content, colors)
    
    # Slide 7: Model Performance
    print("  Adding slide 7: Performance Metrics")
    content = [
        "Outstanding Model Performance:",
        "",
        "Crop Classification:",
        "  ✓ Accuracy: 83.16%",
        "  ✓ Top-3 Accuracy: 99.07%",
        "  ✓ Cross-Validation: 96.28%",
        "",
        "Suitability Prediction:",
        "  ✓ Accuracy: 88.92%",
        "  ✓ ROC AUC: 91.52%",
        "",
        "Yield Prediction:",
        "  ✓ R² Score: 95.96%"
    ]
    add_slide_with_title_and_content(prs, "Model Performance", content, colors)
    
    # Slide 8: Tobacco-Anand Case Study
    print("  Adding slide 8: Case Study")
    content = [
        "Problem Solved: Tobacco-Anand Example",
        "",
        "Original Problem:",
        "  ❌ Tobacco incorrectly predicted as 'Not Grow' in Anand",
        "  ❌ Despite ideal sandy/loamy soil with pH 6.0-7.5",
        "",
        "Solution:",
        "  ✓ Enhanced domain knowledge rules",
        "  ✓ Corrected 3,068 mislabeled records",
        "  ✓ Integrated agricultural expertise",
        "",
        "Result:",
        "  ✅ 100% accuracy on Tobacco-Anand scenarios"
    ]
    add_slide_with_title_and_content(prs, "Problem Solved", content, colors)
    
    # Slide 9: System Workflow
    print("  Adding slide 9: Workflow")
    content = [
        "Complete System Workflow:",
        "",
        "1. User Registration/Login",
        "   → Secure authentication with bcrypt & JWT",
        "",
        "2. Add Soil Details",
        "   → State, District, Taluka, Soil Type, pH",
        "",
        "3. Crop Prediction",
        "   → ML models analyze soil conditions",
        "",
        "4. Results Display",
        "   → Crops, yields, and treatment plans",
        "",
        "5. Translation (Optional)",
        "   → Regional language support"
    ]
    add_slide_with_title_and_content(prs, "System Workflow", content, colors)
    
    # Slide 10: Database Schema
    print("  Adding slide 10: Database")
    content = [
        "PostgreSQL Database:",
        "",
        "1. users table",
        "   • User credentials and profile",
        "",
        "2. soil_details table",
        "   • Soil parameters by location",
        "",
        "3. prediction_history table",
        "   • All predictions with timestamps",
        "   • Confidence scores and yields",
        "   • Full prediction results (JSON)"
    ]
    add_slide_with_title_and_content(prs, "Database Schema", content, colors)
    
    # Slide 11: Security
    print("  Adding slide 11: Security")
    content = [
        "Security Features:",
        "",
        "🔐 Password Security:",
        "   • Bcrypt hashing with salt",
        "",
        "🎫 Session Management:",
        "   • JWT token authentication",
        "",
        "🛡️ Database Security:",
        "   • Parameterized SQL queries",
        "   • SQL injection prevention",
        "",
        "🔒 Data Protection:",
        "   • Environment variables for secrets",
        "   • HTTPS encryption"
    ]
    add_slide_with_title_and_content(prs, "Security & Authentication", content, colors)
    
    # Slide 12: Deployment
    print("  Adding slide 12: Deployment")
    content = [
        "Deployment Architecture:",
        "",
        "Docker Containerization:",
        "  • Python 3.11 slim base image",
        "  • Optimized dependencies",
        "  • Health checks and monitoring",
        "",
        "Railway Platform:",
        "  • Automatic deployment from Git",
        "  • PostgreSQL database provisioning",
        "  • Auto-scaling capabilities"
    ]
    add_slide_with_title_and_content(prs, "Deployment", content, colors)
    
    # Slide 13: Future Enhancements
    print("  Adding slide 13: Future Plans")
    content = [
        "Future Enhancements:",
        "",
        "Short-term:",
        "  • Real-time weather data integration",
        "  • Mobile app development",
        "  • SMS-based recommendations",
        "",
        "Long-term:",
        "  • AI-powered pest detection",
        "  • Satellite imagery integration",
        "  • IoT sensor integration",
        "  • Precision agriculture"
    ]
    add_slide_with_title_and_content(prs, "Future Enhancements", content, colors)
    
    # Slide 14: Impact
    print("  Adding slide 14: Impact")
    content = [
        "Impact & Benefits:",
        "",
        "For Farmers:",
        "  ✓ Data-driven crop decisions",
        "  ✓ Reduced risk of crop failure",
        "  ✓ Increased productivity and income",
        "",
        "For Agriculture:",
        "  ✓ Optimized land utilization",
        "  ✓ Sustainable farming practices",
        "",
        "For Society:",
        "  ✓ Food security improvement",
        "  ✓ Rural economic development"
    ]
    add_slide_with_title_and_content(prs, "Impact & Benefits", content, colors)
    
    # Slide 15: Conclusion
    print("  Adding slide 15: Conclusion")
    content = [
        "Key Achievements:",
        "",
        "✅ 83.16% accuracy with 99.07% top-3 accuracy",
        "✅ Fixed 3,068 mislabeled records",
        "✅ Multi-language support for accessibility",
        "✅ Production-ready with Docker deployment",
        "✅ Real agricultural domain knowledge",
        "",
        "Impact:",
        "Empowering Gujarat farmers with data-driven",
        "agricultural decisions through AI technology"
    ]
    add_slide_with_title_and_content(prs, "Conclusion", content, colors)
    
    # Slide 16: Thank You
    print("  Adding slide 16: Thank You")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.font.size = Pt(60)
    thank_you_para.font.bold = True
    thank_you_para.font.color.rgb = colors['primary']
    thank_you_para.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Questions & Discussion"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(28)
    contact_para.font.color.rgb = colors['text']
    contact_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    filename = 'Gujarat_Crop_Recommendation_Presentation.pptx'
    prs.save(filename)
    
    print("\n" + "="*60)
    print("✅ Presentation created successfully!")
    print("="*60)
    print(f"📄 File: {filename}")
    print(f"📊 Total Slides: {len(prs.slides)}")
    print("\nYou can now open the file with PowerPoint or Google Slides!")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Error: python-pptx library not installed")
        print("📦 Install it using: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
